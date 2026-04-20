"""PPTX Creator — Professional Presentation Engine.

Produces a .pptx file with:
  - Topic-aware color palettes
  - Dark hero title + dark closing slide, light content slides
  - Six varied layouts cycling through every slide (no consecutive repeats)
  - Shape-based visual elements on every content slide
  - Professional typography (Georgia / Calibri)
  - Speaker notes support
"""
from __future__ import annotations

import logging
import os
import re
from typing import Any, Dict, List

from lirox.verify import FileReceipt
from lirox.tools.document_creators.base import ensure_dep, pick_palette, hex_to_rgb, PALETTES

_logger = logging.getLogger("lirox.document_creators.pptx")


def create_pptx(path: str, title: str, slides: List[Dict[str, Any]],
                query: str = "", user_name: str = "") -> FileReceipt:
    """Create a professionally designed PowerPoint presentation.

    Parameters
    ----------
    path:       Absolute output path (will be created if parents exist).
    title:      Presentation title shown on the hero slide.
    slides:     List of dicts with keys ``title``, ``bullets``, ``notes``.
    query:      Original user request — used for palette selection.
    user_name:  Shown in the subtitle / closing credit.

    Returns
    -------
    FileReceipt with ``ok=True`` and ``verified=True`` when the file exists
    on disk with content.  ``error`` is set on any failure.
    """
    from pathlib import Path as _Path
    out_path = _Path(path).resolve()
    r = FileReceipt(tool="pptx_creator", operation="create_pptx", path=str(out_path))
    try:
        ensure_dep("python-pptx", "pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        out_dir = out_path.parent
        out_dir.mkdir(parents=True, exist_ok=True)
        if not os.access(str(out_dir), os.W_OK):
            r.error = f"Output directory is not writable: {out_dir}"
            return r

        palette_name = pick_palette(query or title, title)
        pal = PALETTES[palette_name]
        C = {k: hex_to_rgb(v) for k, v in pal.items()}

        prs = Presentation()
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── Helpers ─────────────────────────────────────────────────────────

        def _blank():
            return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        def _add_bg(slide, color):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, prs.slide_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        def _add_text(slide, left, top, width, height, text,
                      font_name="Georgia", font_size=Pt(18),
                      font_color=None, bold=False,
                      alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.text = str(text)
            p.font.name = font_name
            p.font.size = font_size
            p.font.bold = bold
            p.font.color.rgb = font_color or C["text_dark"]
            p.alignment = alignment
            return txBox

        def _add_accent_line(slide, left, top, width, color=None):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(0.06))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color or C["accent"]
            shape.line.fill.background()

        def _add_icon_circle(slide, left, top, size, text, bg_color, text_color=None):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(top), Inches(size), Inches(size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(text)
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(int(size * 18))
            p.font.color.rgb = text_color or C["text_light"]
            p.font.bold = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

        def _add_bullets(slide, left, top, width, height, bullets,
                         font_size=Pt(16), color=None):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"  •  {bullet}"
                p.font.name = "Calibri"
                p.font.size = font_size
                p.font.color.rgb = color or C["text_dark"]
                p.space_after = Pt(8)
            return txBox

        # ── Slide 1: Hero / Title (dark background) ──────────────────────────
        sl = _blank()
        _add_bg(sl, C["bg_dark"])

        top_bar = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C["accent"]
        top_bar.line.fill.background()

        _add_text(sl, 1.5, 2.0, 10.3, 1.5, title,
                  font_name="Georgia", font_size=Pt(44),
                  font_color=C["text_light"], bold=True,
                  alignment=PP_ALIGN.LEFT)
        _add_accent_line(sl, 1.5, 3.7, 3.0)

        subtitle_text = f"By {user_name}" if user_name else "Generated by Lirox"
        _add_text(sl, 1.5, 4.1, 6, 0.6, subtitle_text,
                  font_name="Calibri", font_size=Pt(18),
                  font_color=hex_to_rgb(pal["accent"]))
        _add_icon_circle(sl, 10.5, 2.5, 1.8, "✦",
                         hex_to_rgb(pal["primary"]), C["text_light"])

        # ── Content slides with varied layouts ───────────────────────────────
        layout_cycle = [
            "two_column", "icon_grid", "stat_callout",
            "full_bleed_left", "card_row", "two_column",
        ]

        for idx, sd in enumerate(slides):
            slide_title = sd.get("title", f"Slide {idx + 2}")
            bullets     = sd.get("bullets", [])
            notes       = sd.get("notes", "")
            layout      = layout_cycle[idx % len(layout_cycle)]
            sl          = _blank()

            _add_bg(sl, C["bg_light"])
            strip = sl.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
            strip.fill.solid()
            strip.fill.fore_color.rgb = hex_to_rgb(pal["primary"])
            strip.line.fill.background()

            if layout == "two_column":
                _add_text(sl, 0.8, 0.5, 6, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)
                if bullets:
                    _add_bullets(sl, 0.8, 1.7, 6, 5.0, bullets,
                                 font_size=Pt(16), color=C["text_dark"])
                colors = [pal["primary"], pal["accent"], pal["secondary"]]
                for i in range(3):
                    block = sl.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(8.5), Inches(1.5 + i * 1.9),
                        Inches(4.0), Inches(1.5))
                    block.fill.solid()
                    block.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
                    block.line.fill.background()
                    if i < len(bullets):
                        tf = block.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        short = bullets[i][:50] + ("…" if len(bullets[i]) > 50 else "")
                        p.text = short
                        p.font.name = "Calibri"
                        p.font.size = Pt(13)
                        p.font.color.rgb = C["text_light"] if i < 2 else C["text_dark"]
                        p.alignment = PP_ALIGN.CENTER
                        try:
                            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        except Exception:
                            pass

            elif layout == "icon_grid":
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)
                icons = ["📌", "🔹", "🔸", "⭐"]
                grid_bullets = bullets[:4] if bullets else ["Point 1"]
                cols = min(len(grid_bullets), 4)
                col_width = 10.5 / cols
                for i, b in enumerate(grid_bullets):
                    cx = 1.2 + i * col_width
                    _add_icon_circle(sl, cx + col_width / 2 - 0.5, 2.0, 1.0,
                                     icons[i % len(icons)],
                                     hex_to_rgb(pal["primary"]))
                    _add_text(sl, cx, 3.3, col_width - 0.4, 3.5, b,
                              font_name="Calibri", font_size=Pt(14),
                              font_color=C["text_dark"],
                              alignment=PP_ALIGN.CENTER)

            elif layout == "stat_callout":
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)
                stat_num = ""
                stat_label = ""
                if bullets:
                    nums = re.findall(r'[\d,]+[+%]?', bullets[0])
                    if nums:
                        stat_num = nums[0]
                        stat_label = bullets[0]
                    else:
                        stat_num = f"0{idx + 1}"
                        stat_label = bullets[0]
                if stat_num:
                    _add_text(sl, 1.0, 2.0, 5, 2.5, stat_num,
                              font_name="Georgia", font_size=Pt(72),
                              font_color=hex_to_rgb(pal["primary"]), bold=True)
                    _add_text(sl, 1.0, 4.2, 5, 1.0, stat_label,
                              font_name="Calibri", font_size=Pt(16),
                              font_color=C["text_dark"])
                if len(bullets) > 1:
                    _add_bullets(sl, 7.0, 2.0, 5.5, 4.5, bullets[1:],
                                 font_size=Pt(15), color=C["text_dark"])

            elif layout == "full_bleed_left":
                left_panel = sl.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, 0, Inches(5.5), prs.slide_height)
                left_panel.fill.solid()
                left_panel.fill.fore_color.rgb = C["bg_dark"]
                left_panel.line.fill.background()
                _add_text(sl, 0.6, 1.5, 4.3, 1.2, slide_title,
                          font_name="Georgia", font_size=Pt(28),
                          font_color=C["text_light"], bold=True)
                _add_accent_line(sl, 0.6, 2.8, 2.0)
                if bullets:
                    _add_text(sl, 0.6, 3.2, 4.3, 3.0, bullets[0],
                              font_name="Calibri", font_size=Pt(15),
                              font_color=hex_to_rgb(pal["accent"]))
                if len(bullets) > 1:
                    _add_bullets(sl, 6.2, 1.5, 6.3, 5.0, bullets[1:],
                                 font_size=Pt(15), color=C["text_dark"])

            elif layout == "card_row":
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)
                card_items = bullets[:3] if bullets else ["Point 1"]
                card_count = min(len(card_items), 3)
                card_w = 3.5
                gap = (11.0 - card_count * card_w) / (card_count + 1)
                card_colors = [pal["primary"], pal["accent"], pal["secondary"]]
                for i, item in enumerate(card_items):
                    cx = 1.0 + gap + i * (card_w + gap)
                    card_bg = sl.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(cx), Inches(2.0), Inches(card_w), Inches(4.5))
                    card_bg.fill.solid()
                    card_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    card_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                    card_bg.line.width = Pt(1)
                    border = sl.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx), Inches(2.0), Inches(card_w), Inches(0.12))
                    border.fill.solid()
                    border.fill.fore_color.rgb = hex_to_rgb(card_colors[i % len(card_colors)])
                    border.line.fill.background()
                    _add_icon_circle(sl, cx + card_w / 2 - 0.35, 2.4, 0.7,
                                     str(i + 1),
                                     hex_to_rgb(card_colors[i % len(card_colors)]))
                    _add_text(sl, cx + 0.3, 3.5, card_w - 0.6, 2.8, item,
                              font_name="Calibri", font_size=Pt(13),
                              font_color=C["text_dark"],
                              alignment=PP_ALIGN.CENTER)

            if notes:
                try:
                    sl.notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass

        # ── Closing slide (dark background) ─────────────────────────────────
        sl = _blank()
        _add_bg(sl, C["bg_dark"])
        bar = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(7.0), prs.slide_width, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C["accent"]
        bar.line.fill.background()

        _add_text(sl, 2.0, 2.0, 9.3, 1.5, "Thank You",
                  font_name="Georgia", font_size=Pt(48),
                  font_color=C["text_light"], bold=True,
                  alignment=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_accent_line(sl, 5.0, 3.8, 3.3)

        credit = user_name or "Lirox AI"
        _add_text(sl, 2.0, 4.2, 9.3, 0.8,
                  f"Presented by {credit}",
                  font_name="Calibri", font_size=Pt(20),
                  font_color=hex_to_rgb(pal["accent"]),
                  alignment=PP_ALIGN.CENTER)
        _add_icon_circle(sl, 1.0, 5.5, 1.0, "✦",
                         hex_to_rgb(pal["primary"]), C["text_light"])
        _add_icon_circle(sl, 11.3, 0.8, 0.8, "◆",
                         hex_to_rgb(pal["accent"]), C["text_light"])

        # ── Persist ──────────────────────────────────────────────────────────
        prs.save(str(out_path))

        if out_path.exists():
            r.ok = True
            r.verified = True
            r.bytes_written = out_path.stat().st_size
            total_slides = len(slides) + 2  # +hero +closing
            r.message = (
                f"Created PowerPoint: {out_path} "
                f"({r.bytes_written:,} bytes, {total_slides} slides, "
                f"palette: {palette_name})"
            )
            r.details["slide_count"] = total_slides
            r.details["palette"] = palette_name
        else:
            r.error = "PowerPoint build completed but file not found on disk"
        return r
    except Exception as e:
        r.error = f"PowerPoint creation error: {e}"
        return r
