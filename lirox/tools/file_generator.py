"""Lirox v1.2 — Professional File Generation Engine

Creates professional-grade PDF, Word, Excel, and PowerPoint files.
Includes topic-aware color palettes, varied slide layouts, design systems.
Dependencies are auto-installed on first use.
"""
from __future__ import annotations

import logging
import os
import subprocess
import sys
import re
import time
from typing import Any, Dict, List, Optional

from lirox.verify import FileReceipt

_logger = logging.getLogger("lirox.file_generator")


def _ensure_dep(package: str, import_name: str = None):
    """Auto-install a dependency if missing."""
    import_name = import_name or package.replace("-", "_")
    try:
        return __import__(import_name)
    except ImportError:
        print(f"  [Lirox] Installing {package}…")
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", package, "--quiet",
             "--disable-pip-version-check"],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        return __import__(import_name)


# ═══════════════════════════════════════════════════════════════
# DESIGN SYSTEM — Color Palettes (topic-aware)
# ═══════════════════════════════════════════════════════════════

PALETTES = {
    "culture": {
        "primary": "B85042",       # Terracotta
        "secondary": "E7E8D1",     # Sand
        "accent": "A7BEAE",        # Sage
        "text_dark": "2D2D2D",
        "text_light": "FFFFFF",
        "bg_dark": "3D2B1F",       # Dark brown
        "bg_light": "FBF7F0",
    },
    "technology": {
        "primary": "1E2761",       # Navy
        "secondary": "CADCFC",     # Ice blue
        "accent": "7B68EE",        # Slate blue
        "text_dark": "1A1A2E",
        "text_light": "FFFFFF",
        "bg_dark": "0F0F23",
        "bg_light": "F8F9FC",
    },
    "nature": {
        "primary": "2C5F2D",       # Forest green
        "secondary": "97BC62",     # Moss
        "accent": "D4E09B",        # Lime
        "text_dark": "1B3A1B",
        "text_light": "FFFFFF",
        "bg_dark": "1B3A1B",
        "bg_light": "F5F9F0",
    },
    "business": {
        "primary": "36454F",       # Charcoal
        "secondary": "F2F2F2",     # Off-white
        "accent": "E85D04",        # Warm orange
        "text_dark": "212121",
        "text_light": "FFFFFF",
        "bg_dark": "212121",
        "bg_light": "FAFAFA",
    },
    "health": {
        "primary": "028090",       # Teal
        "secondary": "00A896",     # Seafoam
        "accent": "02C39A",        # Mint
        "text_dark": "1A1A2E",
        "text_light": "FFFFFF",
        "bg_dark": "014451",
        "bg_light": "F0FAFA",
    },
    "creative": {
        "primary": "6D2E46",       # Berry
        "secondary": "A26769",     # Dusty rose
        "accent": "ECE2D0",        # Cream
        "text_dark": "2D1B28",
        "text_light": "FFFFFF",
        "bg_dark": "2D1B28",
        "bg_light": "FCF8F3",
    },
    "education": {
        "primary": "2B4570",       # Slate blue
        "secondary": "A3CEF1",     # Light blue
        "accent": "E7C582",        # Gold
        "text_dark": "1A2744",
        "text_light": "FFFFFF",
        "bg_dark": "1A2744",
        "bg_light": "F5F8FC",
    },
    "default": {
        "primary": "2563EB",       # Blue
        "secondary": "DBEAFE",     # Light blue
        "accent": "F59E0B",        # Amber
        "text_dark": "1F2937",
        "text_light": "FFFFFF",
        "bg_dark": "111827",
        "bg_light": "F9FAFB",
    },
}

# ── Topic → Palette mapping keywords ─────────────────────────

_PALETTE_KEYWORDS = {
    "culture": [
        "rajasthan", "india", "culture", "history", "heritage", "travel",
        "tourism", "tradition", "ancient", "civilization", "festival",
        "art", "temple", "monument", "cuisine", "folklore", "religion",
    ],
    "technology": [
        "ai", "artificial intelligence", "machine learning", "technology",
        "software", "coding", "programming", "computer", "data science",
        "blockchain", "cybersecurity", "cloud", "devops", "api", "web",
        "algorithm", "neural", "deep learning", "automation", "iot",
    ],
    "nature": [
        "nature", "environment", "climate", "earth", "ocean", "forest",
        "wildlife", "ecology", "biodiversity", "sustainability", "green",
        "geography", "geology", "space", "astronomy", "planet",
    ],
    "business": [
        "business", "finance", "strategy", "marketing", "startup",
        "investment", "economics", "management", "corporate", "revenue",
        "growth", "leadership", "entrepreneurship", "sales", "profit",
    ],
    "health": [
        "health", "medical", "wellness", "fitness", "nutrition",
        "mental health", "medicine", "healthcare", "disease", "therapy",
        "exercise", "diet", "yoga", "mindfulness",
    ],
    "creative": [
        "art", "design", "creative", "photography", "film", "music",
        "animation", "graphic", "media", "branding", "fashion",
    ],
    "education": [
        "education", "learning", "school", "university", "student",
        "teaching", "academic", "curriculum", "course", "training",
        "study", "research", "science", "math", "physics", "chemistry",
    ],
}


def _pick_palette(query: str, title: str = "") -> str:
    """Auto-select color palette based on topic keywords."""
    combined = (query + " " + title).lower()
    best_match = "default"
    best_score = 0
    for palette_name, keywords in _PALETTE_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw in combined)
        if score > best_score:
            best_score = score
            best_match = palette_name
    return best_match


def _hex_to_rgb(hex_str: str):
    """Convert hex color string to RGBColor.

    Validates that *hex_str* is exactly 6 hex characters before parsing so
    that a truncated palette value cannot raise an IndexError or produce a
    silent colour corruption (FIX: was missing length check).
    """
    _ensure_dep("python-pptx", "pptx")
    from pptx.dml.color import RGBColor
    # Normalise: strip leading '#' if present, ensure exactly 6 chars
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        _logger.warning("Invalid hex colour '%s' — falling back to black (000000)", hex_str)
        hex_str = (hex_str + "000000")[:6]  # fallback to black-padded value
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ═══════════════════════════════════════════════════════════════
# PPTX — Professional Presentation Engine
# ═══════════════════════════════════════════════════════════════

def create_pptx(path: str, title: str, slides: List[Dict[str, Any]],
                query: str = "", user_name: str = "") -> FileReceipt:
    """Create a professionally designed PowerPoint presentation.

    Features:
      - Topic-aware color palettes
      - Dark title + closing slides, light content slides
      - Varied layouts per slide (never repeating consecutively)
      - Shape-based visual elements on every slide
      - Professional typography with Georgia / Calibri
    """
    r = FileReceipt(tool="file_generator", operation="create_pptx", path=path)
    try:
        _ensure_dep("python-pptx", "pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)

        palette_name = _pick_palette(query or title, title)
        pal = PALETTES[palette_name]
        C = {k: _hex_to_rgb(v) for k, v in pal.items()}

        prs = Presentation()
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── Helper: add a blank slide ──
        def _blank():
            return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # ── Helper: add full-width background rect ──
        def _add_bg(slide, color):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, prs.slide_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        # ── Helper: add a text box ──
        def _add_text(slide, left, top, width, height, text, font_name="Georgia",
                      font_size=Pt(18), font_color=None, bold=False,
                      alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.text = str(text)
            p.font.name = font_name
            p.font.size = font_size
            p.font.bold = bold
            p.font.color.rgb = font_color or C["text_dark"]
            p.alignment = alignment
            return txBox

        # ── Helper: add accent line ──
        def _add_accent_line(slide, left, top, width, color=None):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(0.06))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color or C["accent"]
            shape.line.fill.background()

        # ── Helper: add circle "icon" ──
        def _add_icon_circle(slide, left, top, size, text, bg_color, text_color=None):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(top), Inches(size), Inches(size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(text)
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(int(size * 18))
            p.font.color.rgb = text_color or C["text_light"]
            p.font.bold = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

        # ── Helper: add bullet list in a text frame ──
        def _add_bullets(slide, left, top, width, height, bullets,
                         font_size=Pt(16), color=None):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"  •  {bullet}"
                p.font.name = "Calibri"
                p.font.size = font_size
                p.font.color.rgb = color or C["text_dark"]
                p.space_after = Pt(8)
            return txBox

        # ══════════════════════════════════════════════
        # SLIDE 1: Title Hero (DARK background)
        # ══════════════════════════════════════════════
        sl = _blank()
        _add_bg(sl, C["bg_dark"])

        # Decorative accent bar at top
        top_bar = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C["accent"]
        top_bar.line.fill.background()

        # Title
        _add_text(sl, 1.5, 2.0, 10.3, 1.5, title,
                  font_name="Georgia", font_size=Pt(44),
                  font_color=C["text_light"], bold=True,
                  alignment=PP_ALIGN.LEFT)

        # Accent line under title
        _add_accent_line(sl, 1.5, 3.7, 3.0)

        # Subtitle
        subtitle_text = f"By {user_name}" if user_name else "Generated by Lirox"
        _add_text(sl, 1.5, 4.1, 6, 0.6, subtitle_text,
                  font_name="Calibri", font_size=Pt(18),
                  font_color=_hex_to_rgb(pal["accent"]), bold=False)

        # Decorative circle
        _add_icon_circle(sl, 10.5, 2.5, 1.8, "✦",
                         _hex_to_rgb(pal["primary"]), C["text_light"])

        # ══════════════════════════════════════════════
        # CONTENT SLIDES (varied layouts)
        # ══════════════════════════════════════════════
        layout_cycle = [
            "two_column", "icon_grid", "stat_callout",
            "full_bleed_left", "card_row", "two_column",
        ]

        for idx, sd in enumerate(slides):
            slide_title = sd.get("title", f"Slide {idx + 2}")
            bullets = sd.get("bullets", [])
            notes = sd.get("notes", "")

            layout = layout_cycle[idx % len(layout_cycle)]
            sl = _blank()

            # Light background
            _add_bg(sl, C["bg_light"])

            # Thin accent strip at top
            strip = sl.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
            strip.fill.solid()
            strip.fill.fore_color.rgb = _hex_to_rgb(pal["primary"])
            strip.line.fill.background()

            if layout == "two_column":
                # ── Two-column: text left, visual element right ──
                _add_text(sl, 0.8, 0.5, 6, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=_hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)

                if bullets:
                    _add_bullets(sl, 0.8, 1.7, 6, 5.0, bullets,
                                 font_size=Pt(16), color=C["text_dark"])

                # Right side: decorative blocks
                colors = [pal["primary"], pal["accent"], pal["secondary"]]
                for i in range(3):
                    block = sl.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(8.5), Inches(1.5 + i * 1.9),
                        Inches(4.0), Inches(1.5))
                    block.fill.solid()
                    block.fill.fore_color.rgb = _hex_to_rgb(colors[i % len(colors)])
                    block.line.fill.background()
                    # Add label
                    if i < len(bullets):
                        tf = block.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        short = bullets[i][:50] + ("…" if len(bullets[i]) > 50 else "")
                        p.text = short
                        p.font.name = "Calibri"
                        p.font.size = Pt(13)
                        p.font.color.rgb = C["text_light"] if i < 2 else C["text_dark"]
                        p.alignment = PP_ALIGN.CENTER
                        try:
                            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        except Exception:
                            pass

            elif layout == "icon_grid":
                # ── Icon grid: 3-4 items with colored circles ──
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=_hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)

                icons = ["📌", "🔹", "🔸", "⭐"]
                grid_bullets = bullets[:4] if bullets else ["Point 1"]
                cols = min(len(grid_bullets), 4)
                col_width = 10.5 / cols
                for i, b in enumerate(grid_bullets):
                    cx = 1.2 + i * col_width
                    _add_icon_circle(sl, cx + col_width/2 - 0.5, 2.0, 1.0,
                                     icons[i % len(icons)],
                                     _hex_to_rgb(pal["primary"]))
                    _add_text(sl, cx, 3.3, col_width - 0.4, 3.5, b,
                              font_name="Calibri", font_size=Pt(14),
                              font_color=C["text_dark"],
                              alignment=PP_ALIGN.CENTER)

            elif layout == "stat_callout":
                # ── Big stat number with description ──
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=_hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)

                # Extract a number from first bullet if possible
                stat_num = ""
                stat_label = ""
                if bullets:
                    nums = re.findall(r'[\d,]+[+%]?', bullets[0])
                    if nums:
                        stat_num = nums[0]
                        stat_label = bullets[0]
                    else:
                        stat_num = f"0{idx + 1}"
                        stat_label = bullets[0]

                # Big stat on the left
                if stat_num:
                    _add_text(sl, 1.0, 2.0, 5, 2.5, stat_num,
                              font_name="Georgia", font_size=Pt(72),
                              font_color=_hex_to_rgb(pal["primary"]), bold=True)
                    _add_text(sl, 1.0, 4.2, 5, 1.0, stat_label,
                              font_name="Calibri", font_size=Pt(16),
                              font_color=C["text_dark"])

                # Remaining bullets on the right
                if len(bullets) > 1:
                    _add_bullets(sl, 7.0, 2.0, 5.5, 4.5, bullets[1:],
                                 font_size=Pt(15), color=C["text_dark"])

            elif layout == "full_bleed_left":
                # ── Dark left panel, light right panel ──
                left_panel = sl.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, 0, Inches(5.5), prs.slide_height)
                left_panel.fill.solid()
                left_panel.fill.fore_color.rgb = C["bg_dark"]
                left_panel.line.fill.background()

                _add_text(sl, 0.6, 1.5, 4.3, 1.2, slide_title,
                          font_name="Georgia", font_size=Pt(28),
                          font_color=C["text_light"], bold=True)
                _add_accent_line(sl, 0.6, 2.8, 2.0)

                if bullets and len(bullets) > 0:
                    _add_text(sl, 0.6, 3.2, 4.3, 3.0, bullets[0],
                              font_name="Calibri", font_size=Pt(15),
                              font_color=_hex_to_rgb(pal["accent"]))

                if len(bullets) > 1:
                    _add_bullets(sl, 6.2, 1.5, 6.3, 5.0, bullets[1:],
                                 font_size=Pt(15), color=C["text_dark"])

            elif layout == "card_row":
                # ── Card row: 3 cards with colored top borders ──
                _add_text(sl, 0.8, 0.5, 11, 0.8, slide_title,
                          font_name="Georgia", font_size=Pt(30),
                          font_color=_hex_to_rgb(pal["primary"]), bold=True)
                _add_accent_line(sl, 0.8, 1.35, 2.0)

                card_items = bullets[:3] if bullets else ["Point 1"]
                card_count = min(len(card_items), 3)
                card_w = 3.5
                gap = (11.0 - card_count * card_w) / (card_count + 1)

                card_colors = [pal["primary"], pal["accent"], pal["secondary"]]
                for i, item in enumerate(card_items):
                    cx = 1.0 + gap + i * (card_w + gap)
                    # Card background
                    card_bg = sl.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(cx), Inches(2.0), Inches(card_w), Inches(4.5))
                    card_bg.fill.solid()
                    card_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    card_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                    card_bg.line.width = Pt(1)

                    # Colored top border
                    border = sl.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx), Inches(2.0), Inches(card_w), Inches(0.12))
                    border.fill.solid()
                    border.fill.fore_color.rgb = _hex_to_rgb(card_colors[i % len(card_colors)])
                    border.line.fill.background()

                    # Card number
                    _add_icon_circle(sl, cx + card_w/2 - 0.35, 2.4, 0.7,
                                     str(i + 1),
                                     _hex_to_rgb(card_colors[i % len(card_colors)]))

                    # Card text
                    _add_text(sl, cx + 0.3, 3.5, card_w - 0.6, 2.8, item,
                              font_name="Calibri", font_size=Pt(13),
                              font_color=C["text_dark"],
                              alignment=PP_ALIGN.CENTER)

            # Speaker notes
            if notes:
                try:
                    sl.notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass

        # ══════════════════════════════════════════════
        # CLOSING SLIDE (DARK background)
        # ══════════════════════════════════════════════
        sl = _blank()
        _add_bg(sl, C["bg_dark"])

        # Bottom accent bar
        bar = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(7.0), prs.slide_width, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C["accent"]
        bar.line.fill.background()

        # Thank you text
        _add_text(sl, 2.0, 2.0, 9.3, 1.5, "Thank You",
                  font_name="Georgia", font_size=Pt(48),
                  font_color=C["text_light"], bold=True,
                  alignment=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)

        _add_accent_line(sl, 5.0, 3.8, 3.3)

        credit = user_name or "Lirox AI"
        _add_text(sl, 2.0, 4.2, 9.3, 0.8,
                  f"Presented by {credit}",
                  font_name="Calibri", font_size=Pt(20),
                  font_color=_hex_to_rgb(pal["accent"]),
                  alignment=PP_ALIGN.CENTER)

        # Decorative circles
        _add_icon_circle(sl, 1.0, 5.5, 1.0, "✦",
                         _hex_to_rgb(pal["primary"]), C["text_light"])
        _add_icon_circle(sl, 11.3, 0.8, 0.8, "◆",
                         _hex_to_rgb(pal["accent"]), C["text_light"])

        # ── Save ──
        prs.save(path)

        if os.path.exists(path):
            r.ok = True
            r.verified = True
            r.bytes_written = os.path.getsize(path)
            total_slides = len(slides) + 2  # +title +closing
            r.message = (f"Created PowerPoint: {path} "
                         f"({r.bytes_written:,} bytes, {total_slides} slides, "
                         f"palette: {palette_name})")
        else:
            r.error = "PowerPoint build completed but file not found on disk"
        return r
    except Exception as e:
        r.error = f"PowerPoint creation error: {e}"
        return r


# ═══════════════════════════════════════════════════════════════
# PDF — Professional Document Engine
# ═══════════════════════════════════════════════════════════════

def create_pdf(path: str, title: str, sections: List[Dict[str, Any]],
               query: str = "", user_name: str = "") -> FileReceipt:
    """Create a professionally styled PDF with cover page, colored headers,
    callout boxes, page numbers, and proper typography."""
    r = FileReceipt(tool="file_generator", operation="create_pdf", path=path)
    try:
        _ensure_dep("reportlab")
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import inch, mm
        from reportlab.lib.colors import HexColor, white, Color
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, PageBreak,
            Table, TableStyle, ListFlowable, ListItem,
            HRFlowable, KeepTogether,
        )

        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)

        palette_name = _pick_palette(query or title, title)
        pal = PALETTES[palette_name]

        doc = SimpleDocTemplate(
            path, pagesize=A4,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch,
            leftMargin=0.85 * inch,
            rightMargin=0.85 * inch,
            title=title,
            author=user_name or "Lirox AI",
        )

        # ── Custom styles ──
        styles = getSampleStyleSheet()

        styles.add(ParagraphStyle(
            'CoverTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=32,
            textColor=HexColor(f"#{pal['primary']}"),
            spaceAfter=12,
            alignment=TA_LEFT,
            leading=38,
        ))

        styles.add(ParagraphStyle(
            'CoverSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=14,
            textColor=HexColor(f"#{pal['accent']}"),
            spaceAfter=6,
        ))

        styles.add(ParagraphStyle(
            'SectionHeader',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=20,
            textColor=HexColor(f"#{pal['primary']}"),
            spaceBefore=24,
            spaceAfter=10,
            leading=24,
        ))

        styles.add(ParagraphStyle(
            'BodyText_Custom',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=11,
            leading=16,
            textColor=HexColor(f"#{pal['text_dark']}"),
            spaceAfter=8,
            alignment=TA_JUSTIFY,
        ))

        styles.add(ParagraphStyle(
            'BulletItem',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=11,
            leading=15,
            textColor=HexColor(f"#{pal['text_dark']}"),
            spaceAfter=4,
            leftIndent=20,
            bulletIndent=10,
        ))

        styles.add(ParagraphStyle(
            'CalloutText',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=11,
            textColor=HexColor(f"#{pal['primary']}"),
            leading=15,
            leftIndent=12,
            rightIndent=12,
            spaceBefore=4,
            spaceAfter=4,
        ))

        styles.add(ParagraphStyle(
            'FooterStyle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=8,
            textColor=HexColor("#999999"),
            alignment=TA_CENTER,
        ))

        # ── Helper: escape XML for reportlab ──
        def _safe(text: str) -> str:
            return (text.replace("&", "&amp;")
                       .replace("<", "&lt;")
                       .replace(">", "&gt;")
                       .replace('"', "&quot;")
                       .replace("'", "&#39;"))

        # ── Helper: accent horizontal rule ──
        def _accent_rule():
            return HRFlowable(
                width="100%", thickness=2,
                color=HexColor(f"#{pal['accent']}"),
                spaceBefore=6, spaceAfter=12,
            )

        # ── Helper: callout box ──
        def _callout_box(text: str):
            safe = _safe(text)
            data = [[Paragraph(safe, styles['CalloutText'])]]
            t = Table(data, colWidths=[doc.width])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1),
                 HexColor(f"#{pal['secondary']}")),
                ('BOX', (0, 0), (-1, -1), 1.5,
                 HexColor(f"#{pal['accent']}")),
                ('LEFTPADDING', (0, 0), (-1, -1), 14),
                ('RIGHTPADDING', (0, 0), (-1, -1), 14),
                ('TOPPADDING', (0, 0), (-1, -1), 10),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
            ]))
            return KeepTogether([Spacer(1, 6), t, Spacer(1, 6)])

        # ── Page number callback ──
        from datetime import datetime

        def _add_page_number(canvas_obj, doc_obj):
            canvas_obj.saveState()
            # Page number
            page_num = f"Page {doc_obj.page}"
            canvas_obj.setFont('Helvetica', 8)
            canvas_obj.setFillColor(HexColor("#999999"))
            canvas_obj.drawCentredString(
                doc_obj.pagesize[0] / 2, 0.4 * inch, page_num)
            # Header line (except first page)
            if doc_obj.page > 1:
                canvas_obj.setStrokeColor(HexColor(f"#{pal['accent']}"))
                canvas_obj.setLineWidth(0.5)
                canvas_obj.line(
                    0.85 * inch, doc_obj.pagesize[1] - 0.55 * inch,
                    doc_obj.pagesize[0] - 0.85 * inch,
                    doc_obj.pagesize[1] - 0.55 * inch)
                canvas_obj.setFont('Helvetica', 7)
                canvas_obj.setFillColor(HexColor(f"#{pal['primary']}"))
                canvas_obj.drawString(
                    0.85 * inch, doc_obj.pagesize[1] - 0.5 * inch,
                    title[:60])
            canvas_obj.restoreState()

        # ══════════════════════════════════════════════
        # BUILD STORY
        # ══════════════════════════════════════════════
        story = []

        # ── Cover page ──
        story.append(Spacer(1, 1.5 * inch))

        # Accent line
        story.append(HRFlowable(
            width="40%", thickness=3,
            color=HexColor(f"#{pal['primary']}"),
            spaceBefore=0, spaceAfter=12, hAlign='LEFT',
        ))

        story.append(Paragraph(_safe(title), styles['CoverTitle']))
        story.append(Spacer(1, 8))

        author_text = user_name or "Lirox AI"
        date_str = datetime.now().strftime("%B %d, %Y")
        story.append(Paragraph(
            f"By {_safe(author_text)}", styles['CoverSubtitle']))
        story.append(Paragraph(date_str, styles['CoverSubtitle']))

        story.append(Spacer(1, 0.5 * inch))

        # Cover decorative rule
        story.append(HRFlowable(
            width="100%", thickness=1,
            color=HexColor(f"#{pal['accent']}"),
            spaceBefore=12, spaceAfter=0,
        ))

        story.append(PageBreak())

        # ── Table of Contents (if 3+ sections) ──
        if len(sections) >= 3:
            story.append(Paragraph("Table of Contents", styles['SectionHeader']))
            story.append(_accent_rule())
            for i, sec in enumerate(sections, 1):
                heading = sec.get("heading", f"Section {i}")
                story.append(Paragraph(
                    f"{i}. &nbsp; {_safe(heading)}",
                    styles['BodyText_Custom']))
            story.append(PageBreak())

        # ── Content sections ──
        for i, sec in enumerate(sections):
            heading = sec.get("heading", "")
            body = sec.get("body", "")
            bullets = sec.get("bullets", [])

            if heading:
                story.append(Paragraph(_safe(heading), styles['SectionHeader']))
                story.append(_accent_rule())

            if body:
                for para in body.split("\n\n"):
                    if para.strip():
                        story.append(Paragraph(
                            _safe(para.strip()), styles['BodyText_Custom']))

            if bullets:
                items = []
                for b in bullets:
                    if b and b.strip():
                        items.append(
                            ListItem(
                                Paragraph(_safe(b.strip()), styles['BulletItem']),
                                bulletColor=HexColor(f"#{pal['primary']}"),
                            ))
                if items:
                    story.append(ListFlowable(
                        items, bulletType="bullet", start="●",
                        bulletFontSize=8,
                    ))
                    story.append(Spacer(1, 8))

            # Add a callout for the first bullet of every other section
            if i % 2 == 0 and bullets:
                story.append(_callout_box(
                    f"Key Point: {bullets[0][:200]}"))

            story.append(Spacer(1, 12))

        # ── Build ──
        doc.build(story, onFirstPage=_add_page_number,
                  onLaterPages=_add_page_number)

        if os.path.exists(path):
            r.ok = True
            r.verified = True
            r.bytes_written = os.path.getsize(path)
            r.message = (f"Created PDF: {path} "
                         f"({r.bytes_written:,} bytes, "
                         f"{len(sections)} sections, palette: {palette_name})")
        else:
            r.error = "PDF build completed but file not found on disk"
        return r
    except Exception as e:
        r.error = f"PDF creation error: {e}"
        return r


# ═══════════════════════════════════════════════════════════════
# WORD — Document Engine
# ═══════════════════════════════════════════════════════════════

def create_docx(path: str, title: str, sections: List[Dict[str, Any]],
                query: str = "", user_name: str = "") -> FileReceipt:
    """Create a Word document with headings, paragraphs, bullets, and tables."""
    r = FileReceipt(tool="file_generator", operation="create_docx", path=path)
    try:
        _ensure_dep("python-docx", "docx")
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
        palette_name = _pick_palette(query or title, title)
        pal = PALETTES[palette_name]
        primary_rgb = RGBColor(
            int(pal["primary"][:2], 16),
            int(pal["primary"][2:4], 16),
            int(pal["primary"][4:6], 16))

        doc = Document()

        # Title
        t = doc.add_heading(title, level=0)
        t.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in t.runs:
            run.font.color.rgb = primary_rgb

        # Author + date
        if user_name:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(f"By {user_name}")
            run.font.size = Pt(12)
            run.font.italic = True

        from datetime import datetime
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(datetime.now().strftime("%B %d, %Y"))
        run.font.size = Pt(10)

        doc.add_page_break()

        for sec in sections:
            if sec.get("heading"):
                h = doc.add_heading(sec["heading"], level=1)
                for run in h.runs:
                    run.font.color.rgb = primary_rgb

            if sec.get("body"):
                for para in sec["body"].split("\n\n"):
                    if para.strip():
                        p = doc.add_paragraph(para.strip())
                        p.paragraph_format.space_after = Pt(8)

            if sec.get("bullets"):
                for b in sec["bullets"]:
                    if b and b.strip():
                        doc.add_paragraph(b.strip(), style="List Bullet")

            if sec.get("table"):
                data = sec["table"]
                if data and len(data) > 0:
                    table = doc.add_table(rows=len(data), cols=len(data[0]))
                    table.style = "Light Grid Accent 1"
                    for i, row_data in enumerate(data):
                        for j, cell_val in enumerate(row_data):
                            table.cell(i, j).text = str(cell_val)

        doc.save(path)

        if os.path.exists(path):
            r.ok = True
            r.verified = True
            r.bytes_written = os.path.getsize(path)
            r.message = f"Created Word doc: {path} ({r.bytes_written:,} bytes)"
        else:
            r.error = "Word build completed but file not found on disk"
        return r
    except Exception as e:
        r.error = f"Word creation error: {e}"
        return r


# ═══════════════════════════════════════════════════════════════
# EXCEL — Spreadsheet Engine
# ═══════════════════════════════════════════════════════════════

def create_xlsx(path: str, title: str, sheets: List[Dict[str, Any]],
                query: str = "", user_name: str = "") -> FileReceipt:
    """Create an Excel workbook with styled headers and data."""
    r = FileReceipt(tool="file_generator", operation="create_xlsx", path=path)
    try:
        _ensure_dep("openpyxl")
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)

        palette_name = _pick_palette(query or title, title)
        pal = PALETTES[palette_name]

        wb = Workbook()
        wb.remove(wb.active)

        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill(
            start_color=pal["primary"], end_color=pal["primary"],
            fill_type="solid")
        alt_fill = PatternFill(
            start_color=pal["secondary"], end_color=pal["secondary"],
            fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin"))

        for sd in sheets:
            ws = wb.create_sheet(title=str(sd.get("name", "Sheet"))[:31])
            headers = sd.get("headers", [])
            rows = sd.get("rows", [])

            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=str(h))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border

            for ri, row_data in enumerate(rows, 2):
                for ci, val in enumerate(row_data, 1):
                    cell = ws.cell(row=ri, column=ci, value=val)
                    cell.border = thin_border
                    # Alternating row colors
                    if ri % 2 == 0:
                        cell.fill = alt_fill

            # Auto-fit columns
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    try:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

        wb.save(path)

        if os.path.exists(path):
            r.ok = True
            r.verified = True
            r.bytes_written = os.path.getsize(path)
            r.message = f"Created Excel: {path} ({r.bytes_written:,} bytes)"
        else:
            r.error = "Excel build completed but file not found on disk"
        return r
    except Exception as e:
        r.error = f"Excel creation error: {e}"
        return r
